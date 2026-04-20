"""
File readers for all supported document types.
Each reader returns a list of Section objects — one per natural unit
(page, slide, chapter, heading-section, etc.).
"""

import re
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Section:
    text: str
    phys_page: int  # 1-based index within the document
    page_label: str  # human-readable label used in citations
    total: int  # total sections in this document
    links_to: list[str] = field(default_factory=list)  # basenames of linked HTML files
    atomic: bool = False  # if True: store entire section as one chunk, never split
    breadcrumb: str = ""  # hierarchical heading path, e.g. "Namespace > ClassName"
    cross_refs: list[str] = field(default_factory=list)  # @see / related symbol names


def _make_sections(parts: list[tuple[str, str]]) -> list[Section]:
    """Build Section list from (text, label) pairs."""
    total = len(parts)
    return [
        Section(text=t, phys_page=i + 1, page_label=label, total=total)
        for i, (t, label) in enumerate(parts)
    ]


# ── PDF ───────────────────────────────────────────────────────────────────────


def read_pdf(path: Path) -> list[Section]:
    import fitz  # pymupdf

    doc = fitz.open(str(path))
    total = len(doc)
    sections = []
    for page_num in range(total):
        page = doc[page_num]
        text = page.get_text().strip()
        phys_page = page_num + 1
        label = page.get_label() or str(phys_page)
        sections.append(Section(text=text, phys_page=phys_page, page_label=label, total=total))
    doc.close()
    return sections


# ── Plain text ────────────────────────────────────────────────────────────────


def read_text(path: Path) -> list[Section]:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return [Section(text=text, phys_page=1, page_label="1", total=1)]


# ── Markdown — split by ATX headings (#, ##, …) ──────────────────────────────


def read_markdown(path: Path) -> list[Section]:
    text = path.read_text(encoding="utf-8", errors="replace")
    heading_re = re.compile(r"^#{1,6}\s+(.+)$", re.MULTILINE)
    parts = _split_by_headings(text, heading_re)
    return _make_sections(parts) if parts else read_text(path)


# ── reStructuredText — split by underlined section titles ────────────────────


def read_rst(path: Path) -> list[Section]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    parts: list[tuple[str, str]] = []
    current: list[str] = []
    current_label = "intro"
    underline_re = re.compile(r"^[=\-~^\"'`#*+]{2,}$")
    i = 0
    while i < len(lines):
        line = lines[i]
        next_line = lines[i + 1] if i + 1 < len(lines) else ""
        is_underline = next_line and underline_re.match(next_line) and len(next_line) >= len(line)
        if line.strip() and is_underline:
            if current:
                parts.append(("\n".join(current).strip(), current_label))
            current_label = line.strip()[:60]
            current = [line, next_line]
            i += 2
        else:
            current.append(line)
            i += 1
    if current:
        parts.append(("\n".join(current).strip(), current_label))
    parts = [(t, lbl) for t, lbl in parts if t.strip()]
    return _make_sections(parts) if parts else read_text(path)


def _split_by_headings(text: str, heading_re: re.Pattern) -> list[tuple[str, str]]:
    parts: list[tuple[str, str]] = []
    current: list[str] = []
    current_label = "intro"
    for line in text.splitlines():
        m = heading_re.match(line)
        if m:
            if current:
                parts.append(("\n".join(current).strip(), current_label))
            current_label = m.group(1).strip()[:60]
            current = [line]
        else:
            current.append(line)
    if current:
        parts.append(("\n".join(current).strip(), current_label))
    return [(t, lbl) for t, lbl in parts if t.strip()]


# ── HTML ──────────────────────────────────────────────────────────────────────


def _extract_html_links(soup, source_name: str | None = None) -> list[str]:
    """Return sorted list of relative HTML filenames linked from *soup*.

    Skips external URLs, anchor-only hrefs, non-HTML targets, and self-links.
    Returns basenames only so they match how filenames are stored in the index.
    """
    from urllib.parse import urlparse

    seen: set[str] = set()
    for tag in soup.find_all("a", href=True):
        href = tag["href"].strip()
        if not href:
            continue
        parsed = urlparse(href)
        if parsed.scheme or parsed.netloc:
            continue  # external URL
        path_part = parsed.path
        if not path_part:
            continue  # anchor-only (#fragment)
        if not (path_part.endswith(".html") or path_part.endswith(".htm")):
            continue
        basename = Path(path_part).name
        if not basename or basename == source_name:
            continue  # self-link
        seen.add(basename)
    return sorted(seen)


def _process_section_nodes(nodes: list, heading_level: int) -> tuple[str, list[str], bool]:
    """Extract text, cross-references, and atomic status from a list of BeautifulSoup nodes.

    Returns (text, cross_refs, is_atomic).
    is_atomic is True for h4-level sections (single-declaration in Doxygen) and for
    sections containing a Doxygen enum fieldtable or dl.enum element.
    Cross-refs are extracted from <dl class="section see"> blocks and removed from text.
    """
    from bs4 import NavigableString, Tag

    lines: list[str] = []
    cross_refs: list[str] = []
    is_atomic = heading_level >= 4  # h4 = single-declaration in Doxygen

    for node in nodes:
        if isinstance(node, Tag):
            node_classes = set(node.get("class") or [])
            # If this node itself is a top-level @see block, extract refs and skip
            if node.name == "dl" and {"section", "see"} <= node_classes:
                for item in node.find_all(["a", "dd"]):
                    ref = item.get_text(" ", strip=True)
                    if ref:
                        cross_refs.append(ref)
                continue  # do not include @see text in section body
            # Extract nested @see cross-refs, remove them before text flattening
            for dl in node.find_all("dl", class_="section see"):
                for item in dl.find_all(["a", "dd"]):
                    ref = item.get_text(" ", strip=True)
                    if ref:
                        cross_refs.append(ref)
                dl.decompose()
            # Doxygen enum fieldtable or explicit enum dl → treat section as atomic
            # Check both the node itself and any nested matches
            is_fieldtable = node.name == "table" and "fieldtable" in node_classes
            is_enum_dl = node.name == "dl" and "enum" in node_classes
            if (
                is_fieldtable
                or is_enum_dl
                or node.find("table", class_="fieldtable")
                or node.find("dl", class_="enum")
            ):
                is_atomic = True
            # Flatten remaining text
            for child in node.descendants:
                if isinstance(child, NavigableString):
                    line = str(child).strip()
                    if line:
                        lines.append(line)
        elif isinstance(node, NavigableString):
            line = str(node).strip()
            if line:
                lines.append(line)

    return "\n".join(lines).strip(), sorted(set(cross_refs)), is_atomic


def read_html(path: Path) -> list[Section]:
    from bs4 import BeautifulSoup, NavigableString, Tag

    html = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")

    links = _extract_html_links(soup, source_name=path.name)

    for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
        tag.decompose()
    for tag in soup("ul", class_="ccore"):  # STW core-availability badges (sc/bc/ac)
        tag.decompose()

    heading_tags = {"h1", "h2", "h3", "h4"}
    heading_level_map = {"h1": 1, "h2": 2, "h3": 3, "h4": 4}
    default_label = (
        soup.title.string.strip()[:60] if soup.title and soup.title.string else path.stem
    )

    # parts: (text, label, breadcrumb, atomic, cross_refs)
    parts: list[tuple[str, str, str, bool, list[str]]] = []
    current_nodes: list = []
    current_label = default_label
    current_level = 0
    heading_stack: list[tuple[int, str]] = []  # (level, label)

    # Prefer a semantic content container so deeply-nested Sphinx/Bootstrap layouts
    # don't fall through to the descendants-based fallback with boilerplate noise.
    body = soup.body or soup
    root = (
        body.find("main")
        or body.find(attrs={"role": "main"})
        or body.find("article")
        or body.find("div", class_="rst-content")
        or body
    )
    for node in root.children:
        if isinstance(node, Tag) and node.name in heading_tags:
            # Flush accumulated section
            text, xrefs, atomic = _process_section_nodes(current_nodes, current_level)
            if text.strip():
                bc = " > ".join(lbl for _, lbl in heading_stack[:-1])
                parts.append((text, current_label, bc, atomic, xrefs))
            # Update heading stack: pop levels >= current heading level
            level = heading_level_map[node.name]
            heading_stack = [(lvl, lbl) for lvl, lbl in heading_stack if lvl < level]
            new_label = node.get_text(" ", strip=True)[:60] or current_label
            heading_stack.append((level, new_label))
            current_label, current_level, current_nodes = new_label, level, []
        else:
            current_nodes.append(node)

    # Flush final section
    text, xrefs, atomic = _process_section_nodes(current_nodes, current_level)
    if text.strip():
        bc = " > ".join(lbl for _, lbl in heading_stack[:-1])
        parts.append((text, current_label, bc, atomic, xrefs))

    parts = [(t, lbl, bc, at, xr) for t, lbl, bc, at, xr in parts if t.strip()]

    # Fallback: if top-level children traversal found nothing (headings nested in divs),
    # fall back to the descendants-based approach without the new features.
    if not parts:
        # Re-parse from scratch since we may have mutated the soup above
        soup2 = BeautifulSoup(html, "html.parser")
        for tag in soup2(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        fallback_parts: list[tuple[str, str]] = []
        fb_lines: list[str] = []
        fb_label = default_label
        for node in (soup2.body or soup2).descendants:
            if isinstance(node, Tag) and node.name in heading_tags:
                fb_text = "\n".join(fb_lines).strip()
                if fb_text:
                    fallback_parts.append((fb_text, fb_label))
                fb_label = node.get_text(" ", strip=True)[:60] or fb_label
                fb_lines = []
            elif isinstance(node, NavigableString) and not any(
                isinstance(p, Tag) and p.name in heading_tags for p in node.parents
            ):
                line = str(node).strip()
                if line:
                    fb_lines.append(line)
        fb_text = "\n".join(fb_lines).strip()
        if fb_text:
            fallback_parts.append((fb_text, fb_label))
        fallback_parts = [(t, lbl) for t, lbl in fallback_parts if t.strip()]
        sections = _make_sections(fallback_parts) if fallback_parts else read_text(path)
        if links:
            for s in sections:
                s.links_to = links
        return sections

    total = len(parts)
    sections = [
        Section(
            text=t,
            phys_page=i + 1,
            page_label=lbl,
            total=total,
            links_to=links,
            breadcrumb=bc,
            atomic=at,
            cross_refs=xr,
        )
        for i, (t, lbl, bc, at, xr) in enumerate(parts)
    ]
    return sections


# ── DOCX ──────────────────────────────────────────────────────────────────────


def read_docx(path: Path) -> list[Section]:
    from docx import Document

    doc = Document(str(path))
    parts: list[tuple[str, str]] = []
    current: list[str] = []
    current_label = "intro"

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading") and para.text.strip():
            if current:
                parts.append(("\n".join(current).strip(), current_label))
            current_label = para.text.strip()[:60]
            current = []
        elif para.text.strip():
            current.append(para.text.strip())

    if current:
        parts.append(("\n".join(current).strip(), current_label))

    parts = [(t, lbl) for t, lbl in parts if t.strip()]
    return _make_sections(parts) if parts else read_text(path)


# ── PPTX ──────────────────────────────────────────────────────────────────────


def read_pptx(path: Path) -> list[Section]:
    from pptx import Presentation

    prs = Presentation(str(path))
    total = len(prs.slides)
    sections = []

    for i, slide in enumerate(prs.slides):
        lines: list[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    lines.append(line)
            if not title and shape.text.strip():
                title = shape.text.strip()[:60]

        text = "\n".join(lines).strip()
        label = title or f"Slide {i + 1}"
        sections.append(Section(text=text, phys_page=i + 1, page_label=label, total=total))

    return sections


# ── EPUB ──────────────────────────────────────────────────────────────────────


def read_epub(path: Path) -> list[Section]:
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(str(path), options={"ignore_ncx": True})
    parts: list[tuple[str, str]] = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n").strip()
        if not text:
            continue
        heading = soup.find(["h1", "h2", "h3"])
        label = heading.get_text().strip()[:60] if heading else item.get_name()
        parts.append((text, label))

    return _make_sections(parts) if parts else read_text(path)


# ── Registry ──────────────────────────────────────────────────────────────────

READERS: dict[str, Callable[[Path], list[Section]]] = {
    ".pdf": read_pdf,
    ".txt": read_text,
    ".md": read_markdown,
    ".rst": read_rst,
    ".html": read_html,
    ".htm": read_html,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".epub": read_epub,
}

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(READERS)


def read_file(path: Path) -> list[Section]:
    ext = path.suffix.lower()
    if ext not in READERS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type: {ext}. Supported: {supported}")
    return READERS[ext](path)
